from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 0 = Title Slide
    # 1 = Title and Content
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "AgenticChain Finance"
    subtitle1.text = "Transforming Supply Chain Finance with Agentic AI & Blockchain\n\nProof of Concept MVP"
    
    # Slide 2: The Problem & Solution
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = "Bridging the Trust Gap"
    tf = slide2.placeholders[1].text_frame
    tf.text = "The Trust Gap:"
    p = tf.add_paragraph()
    p.text = "Banks lack real-time visibility into SME manufacturing operations, increasing risk."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "SMEs struggle to access Just-In-Time (JIT) financing without verifiable physicals."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Our Solution: AgenticChain"
    p = tf.add_paragraph()
    p.text = "An AI Agent verifies production physicals (IoT + ERP) in real-time."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Verified milestones are secured immutably on a Blockchain ledger."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Unlocks dynamic Smart Contract financing for Bankers based on verified truth."
    p.level = 1
    
    img_arch = "/Users/user/.gemini/antigravity/brain/50714f60-3187-435b-a57e-3c5d1c90831a/system_architecture_1773569338530.png"
    img_sme = "/Users/user/.gemini/antigravity/brain/50714f60-3187-435b-a57e-3c5d1c90831a/sme_operations_1773569276480.png"
    img_ai = "/Users/user/.gemini/antigravity/brain/50714f60-3187-435b-a57e-3c5d1c90831a/ai_auditor_1773569293063.png"
    img_chain = "/Users/user/.gemini/antigravity/brain/50714f60-3187-435b-a57e-3c5d1c90831a/blockchain_ledger_1773569307080.png"
    img_bank = "/Users/user/.gemini/antigravity/brain/50714f60-3187-435b-a57e-3c5d1c90831a/banker_portal_1773569323261.png"

    # Slide 3: Core Architecture & Personas
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    slide3.shapes.title.text = "System Architecture & 4 Personas"
    slide3.placeholders[1].width = Inches(5.0)
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "A holistic trust loop simulation comprising four core modules:"
    p = tf3.add_paragraph()
    p.text = "1. SME Operations: The physical layer (ERP & IoT streams)."
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "2. AI Auditor: The real-time reasoning and verification engine."
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "3. Blockchain Ledger: The immutable state anchor for confirmed physicals."
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "4. Banker's Portal: The financial interface for dynamic approvals."
    p.level = 1
    slide3.shapes.add_picture(img_arch, Inches(5.5), Inches(2.0), width=Inches(4.0))
    
    # Slide 4: Persona 1 - SME Operations
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    slide4.shapes.title.text = "🏭 Persona 1: SME Operations Dashboard"
    slide4.placeholders[1].width = Inches(5.0)
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Role: SME Manager"
    p = tf4.add_paragraph()
    p.text = "Synchronize ERP Data: Upload Purchase Orders (PO) and Raw Material Inventory."
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Connect IoT Factory Feeds: Stream Machine Uptime and Units Produced."
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Run Productions: Simulates either a Success scenario or a Risk (Anomaly) scenario."
    p.level = 1
    slide4.shapes.add_picture(img_sme, Inches(5.5), Inches(2.0), width=Inches(4.0))
    
    # Slide 5: Persona 2 - AI Auditor
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    slide5.shapes.title.text = "🤖 Persona 2: Agentic Reasoning Engine"
    slide5.placeholders[1].width = Inches(5.0)
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Role: Automated Real-Time Auditor"
    p = tf5.add_paragraph()
    p.text = "Cross-Verification: Compares raw material inventory against PO requirements."
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Anomaly Detection: Monitors IoT metrics (e.g., machine downtime below 90%)."
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Action Engine: Automatically boosts credit score if nominal, or alerts Banker and drops score if anomalies are found."
    p.level = 1
    slide5.shapes.add_picture(img_ai, Inches(5.5), Inches(2.0), width=Inches(4.0))
    
    # Slide 6: Persona 3 - Blockchain Trust Layer
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    slide6.shapes.title.text = "🔗 Persona 3: Blockchain Ledger"
    slide6.placeholders[1].width = Inches(5.0)
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Role: Immutable Record of Truth"
    p = tf6.add_paragraph()
    p.text = "Cryptographic Hashes: The AI Auditor generates an immutable SHA-256 state anchor for successful production runs."
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "Transparency: Provides mathematical certainty of physical operations."
    p.level = 1
    p = tf6.add_paragraph()
    p.text = "Smart Contract Fuel: Serves as the trusted trigger for financial execution."
    p.level = 1
    slide6.shapes.add_picture(img_chain, Inches(5.5), Inches(2.0), width=Inches(4.0))
    
    # Slide 7: Persona 4 - Banker's Portal
    slide7 = prs.slides.add_slide(bullet_slide_layout)
    slide7.shapes.title.text = "🏦 Persona 4: Banker's Portal"
    slide7.placeholders[1].width = Inches(5.0)
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Role: Bank Officer / Automated Financier"
    p = tf7.add_paragraph()
    p.text = "Live Credit Score: A dynamic gauge that updates purely on verified physicals."
    p.level = 1
    p = tf7.add_paragraph()
    p.text = "Automated Workflows: Only allows funds release via Smart Contract if:"
    p.level = 1
    p = tf7.add_paragraph()
    p.text = "Score > 750"
    p.level = 2
    p = tf7.add_paragraph()
    p.text = "No AI-detected anomalies exist"
    p.level = 2
    p = tf7.add_paragraph()
    p.text = "At least one ledger anchor is verified"
    p.level = 2
    slide7.shapes.add_picture(img_bank, Inches(5.5), Inches(2.0), width=Inches(4.0))
    
    # Slide 8: Conclusion
    slide8 = prs.slides.add_slide(bullet_slide_layout)
    slide8.shapes.title.text = "Outcome & Value Proposition"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "For SMEs:"
    p = tf8.add_paragraph()
    p.text = "Access cheaper, faster capital securely."
    p.level = 1
    p = tf8.add_paragraph()
    p.text = "For Banks:"
    p = tf8.add_paragraph()
    p.text = "De-risk supply chain lending with verifiable, real-time ground truth."
    p.level = 1
    p = tf8.add_paragraph()
    p.text = "The Future of Industrial Fintech:"
    p = tf8.add_paragraph()
    p.text = "Merging industrial physicals with decentralized logic."
    p.level = 1

    ppt_path = "AgenticChain_Finance_MVP_Presentation_Visuals.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved to {ppt_path}")

if __name__ == "__main__":
    create_presentation()
